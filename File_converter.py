import os
from docx2pdf import convert
import pdfkit
from fpdf import FPDF
from PIL import Image
import pdf2image
from pptx import Presentation
import pandas as pd

# Conversion functions
def word_to_pdf(docx_path):
    """Convert Word document to PDF using docx2pdf."""
    pdf_path = os.path.splitext(docx_path)[0] + ".pdf"
    try:
        convert(docx_path, pdf_path)
        print(f"Converted {docx_path} to {pdf_path}")
    except Exception as e:
        print(f"Error converting Word to PDF: {e}")

def txt_to_pdf(txt_path):
    """Convert text file to PDF."""
    pdf_path = os.path.splitext(txt_path)[0] + ".pdf"
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    try:
        with open(txt_path, 'r') as file:
            for line in file:
                pdf.cell(200, 10, txt=line.encode('latin-1', 'replace').decode('latin-1'), ln=True)
        pdf.output(pdf_path)
        print(f"Converted {txt_path} to {pdf_path}")
    except Exception as e:
        print(f"Error converting TXT to PDF: {e}")

def png_to_jpeg(png_path):
    """Convert PNG image to JPEG."""
    jpeg_path = os.path.splitext(png_path)[0] + ".jpg"
    try:
        image = Image.open(png_path)
        rgb_im = image.convert('RGB')
        rgb_im.save(jpeg_path, format="JPEG")
        print(f"Converted {png_path} to {jpeg_path}")
    except Exception as e:
        print(f"Error converting PNG to JPEG: {e}")

def pdf_to_pptx(pdf_path):
    """Convert PDF to PPTX slides."""
    pptx_path = os.path.splitext(pdf_path)[0] + ".pptx"
    try:
        images = pdf2image.convert_from_path(pdf_path)
        prs = Presentation()
        for image in images:
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            img_path = "temp_image.jpg"
            image.save(img_path, 'JPEG')
            slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)
            os.remove(img_path)
        prs.save(pptx_path)
        print(f"Converted {pdf_path} to {pptx_path}")
    except Exception as e:
        print(f"Error converting PDF to PPTX: {e}")

def excel_to_pdf(excel_path):
    """Convert Excel file to PDF."""
    pdf_path = os.path.splitext(excel_path)[0] + ".pdf"
    try:
        df = pd.read_excel(excel_path)
        temp_html = "temp.html"
        df.to_html(temp_html)
        pdfkit.from_file(temp_html, pdf_path)
        os.remove(temp_html)
        print(f"Converted {excel_path} to {pdf_path}")
    except Exception as e:
        print(f"Error converting Excel to PDF: {e}")

def html_to_pdf(html_path):
    """Convert HTML file to PDF."""
    pdf_path = os.path.splitext(html_path)[0] + ".pdf"
    try:
        pdfkit.from_file(html_path, pdf_path)
        print(f"Converted {html_path} to {pdf_path}")
    except Exception as e:
        print(f"Error converting HTML to PDF: {e}")

def jpeg_to_png(jpeg_path):
    """Convert JPEG image to PNG."""
    png_path = os.path.splitext(jpeg_path)[0] + ".png"
    try:
        image = Image.open(jpeg_path)
        image.save(png_path, format="PNG")
        print(f"Converted {jpeg_path} to {png_path}")
    except Exception as e:
        print(f"Error converting JPEG to PNG: {e}")

# Main function for auto-detection
def convert_file(file_path):
    extension = os.path.splitext(file_path)[1].lower()
    
    if extension == ".docx":
        word_to_pdf(file_path)
    elif extension == ".txt":
        txt_to_pdf(file_path)
    elif extension == ".png":
        png_to_jpeg(file_path)
    elif extension == ".pdf":
        # Assume conversion to PPTX as an example; other conversions could be added
        pdf_to_pptx(file_path)
    elif extension == ".xlsx":
        excel_to_pdf(file_path)
    elif extension == ".html":
        html_to_pdf(file_path)
    elif extension == ".jpeg" or extension == ".jpg":
        jpeg_to_png(file_path)
    else:
        print(f"Unsupported file type: {extension}")

# Example usage
if __name__ == "__main__":
    file_path = input("Enter the path of the file you want to convert: ")
    convert_file(file_path)
